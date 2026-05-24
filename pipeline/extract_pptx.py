from __future__ import annotations

from pathlib import Path

from pipeline.models import Block


def extract_pptx(path: Path) -> list[Block]:
    from pptx import Presentation

    prs = Presentation(path)
    blocks: list[Block] = []
    doc_id = path.stem

    for slide_index, slide in enumerate(prs.slides, start=1):
        title = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text else f"Slide {slide_index}"
        title_bbox = _shape_bbox(slide.shapes.title) if slide.shapes.title else None
        blocks.append(
            Block(
                block_id=f"{doc_id}_s{slide_index}_title",
                doc_id=doc_id,
                source_file=path.name,
                source_type="pptx",
                content_type="slide_title",
                text=title,
                slide=slide_index,
                section=title,
                bbox=title_bbox,
                reading_order=0,
                metadata={"extraction_confidence": 0.96},
            )
        )

        order = 1
        for shape_index, shape in enumerate(slide.shapes):
            if shape == slide.shapes.title or not getattr(shape, "has_text_frame", False):
                continue
            text = "\n".join(
                paragraph.text.strip()
                for paragraph in shape.text_frame.paragraphs
                if paragraph.text and paragraph.text.strip()
            )
            if not text:
                continue
            content_type = "speaker_notes" if text.lower().startswith("speaker notes:") else "bullet_list"
            blocks.append(
                Block(
                    block_id=f"{doc_id}_s{slide_index}_shape_{shape_index}",
                    doc_id=doc_id,
                    source_file=path.name,
                    source_type="pptx",
                    content_type=content_type,
                    text=text,
                    slide=slide_index,
                    section=title,
                    bbox=_shape_bbox(shape),
                    reading_order=order,
                    metadata={"extraction_confidence": 0.91},
                )
            )
            order += 1

    return blocks


def _shape_bbox(shape) -> dict[str, float]:
    return {
        "x": float(shape.left),
        "y": float(shape.top),
        "width": float(shape.width),
        "height": float(shape.height),
    }
