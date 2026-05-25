from __future__ import annotations

from pathlib import Path


RAW_DIR = Path("data/raw")


def generate_pdf(path: Path = RAW_DIR / "enterprise_report.pdf") -> Path:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

    path.parent.mkdir(parents=True, exist_ok=True)
    styles = getSampleStyleSheet()
    doc = SimpleDocTemplate(str(path), pagesize=letter, title="Enterprise Report")
    story = []

    def header_footer(canvas, document):
        canvas.saveState()
        canvas.setFont("Helvetica", 9)
        canvas.drawString(54, 760, "Synthetic Enterprise Quarterly Operating Report - Public Demo Data")
        canvas.drawRightString(558, 32, f"Page {document.page}")
        canvas.restoreState()

    pages = [
        (
            "Executive Summary",
            [
                "Acme delivered steady enterprise expansion in Q2. APAC and EMEA grew despite procurement delays.",
                "Customer churn was elevated in the SMB segment, while enterprise renewal rates improved.",
                "Noisy OCR-like text: cust0mer retent1on remalns a board-level metr1c.",
            ],
            None,
        ),
        (
            "Financial Performance",
            [
                "Revenue performance was supported by platform attach rates and support expansion.",
                "The following table begins the quarterly revenue view and continues on the next page.",
            ],
            [
                ["Region", "Revenue", "Growth"],
                ["APAC", "$2.1M", "12%"],
                ["EMEA", "$1.4M", "8%"],
            ],
        ),
        (
            "Quarterly Revenue",
            [
                "The quarterly revenue table continuation shows North America and Latin America.",
                "APAC revenue growth remained the highest among major operating regions.",
            ],
            [
                ["Region", "Revenue", "Growth"],
                ["North America", "$3.2M", "10%"],
                ["LATAM", "$0.7M", "5%"],
            ],
        ),
        (
            "Customer Health",
            [
                "Customer churn appeared in self-service accounts after onboarding friction increased.",
                "Enterprise accounts cited stronger support response and clearer migration paths.",
            ],
            None,
        ),
        (
            "Operational Risks",
            [
                "Primary risks include delayed security reviews, integration capacity, and data residency requirements.",
                "Mitigations include earlier legal intake, standard integration templates, and regional processing controls.",
            ],
            None,
        ),
    ]

    for index, (section, paragraphs, table_rows) in enumerate(pages):
        story.append(Paragraph(section, styles["Heading1"]))
        story.append(Spacer(1, 12))
        for paragraph in paragraphs:
            story.append(Paragraph(paragraph, styles["BodyText"]))
            story.append(Spacer(1, 10))
        if table_rows:
            table = Table(table_rows, hAlign="LEFT")
            table.setStyle(
                TableStyle(
                    [
                        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#D9EAF7")),
                        ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                        ("ALIGN", (1, 1), (-1, -1), "RIGHT"),
                        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ]
                )
            )
            story.append(table)
        if index < len(pages) - 1:
            story.append(PageBreak())

    doc.build(story, onFirstPage=header_footer, onLaterPages=header_footer)
    return path


def generate_pptx(path: Path = RAW_DIR / "strategy_deck.pptx") -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()

    slides = [
        (
            "Enterprise Strategy Overview",
            ["Expand table-aware retrieval", "Launch MCP knowledge tools", "Use citations for audit workflows"],
            "Speaker notes: emphasize traceability and source-grounded answers.",
        ),
        (
            "Risks and Mitigations",
            ["Risk: stale source documents", "Risk: poor OCR quality", "Mitigation: diagnostics and confidence flags"],
            "Speaker notes: mention that quality diagnostics drive review queues.",
        ),
        (
            "MCP Architecture",
            ["LLM client connects to MCP server", "MCP tools retrieve vector chunks", "Resources expose outlines and tables"],
            "Speaker notes: architecture diagram text: Client -> MCP Server -> Vector DB.",
        ),
        (
            "Rollout Metrics",
            ["Target latency: under 800ms", "Evaluation set: 5 source-grounded questions", "Customer churn query must cite report"],
            "Speaker notes: compare semantic retrieval with metadata lookup.",
        ),
    ]

    for title, bullets, notes in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        body = slide.placeholders[1].text_frame
        body.clear()
        for idx, bullet in enumerate(bullets):
            paragraph = body.paragraphs[0] if idx == 0 else body.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
        left, top, width, height = Inches(1), Inches(4.8), Inches(7), Inches(0.6)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        textbox.text_frame.text = notes

    prs.save(path)
    return path


def main() -> None:
    print(generate_pdf())
    print(generate_pptx())


if __name__ == "__main__":
    main()
